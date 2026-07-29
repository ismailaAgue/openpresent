import io
from backend.ports.export import ExportPort
from backend.models.recipe import Recipe, BlockType

_COLOR_SETS = {
    "neutral": {"title": (0x22, 0x22, 0x22), "accent": (0x44, 0x44, 0x44)},
    "blue_academic": {"title": (0x1B, 0x3A, 0x5C), "accent": (0x2E, 0x5C, 0x8A)},
}


class PptxExportAdapter(ExportPort):
    def format_id(self) -> str:
        return "pptx"

    def export(self, recipe: Recipe) -> bytes:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
        except ImportError as e:
            raise RuntimeError(
                "python-pptx is required for PPTX export. Install with: "
                "pip install python-pptx --break-system-packages"
            ) from e

        colors = _COLOR_SETS.get(recipe.theme.color_set_id, _COLOR_SETS["neutral"])
        title_color = RGBColor(*colors["title"])

        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

        for i, slide_data in enumerate(sorted(recipe.outline.slides, key=lambda s: s.order)):
            is_title_slide = (i == 0)
            layout = title_layout if is_title_slide else content_layout
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_data.title
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color

            if is_title_slide:
                continue

            body_texts = [
                b.text for b in slide_data.content_blocks
                if b.type in (BlockType.BULLET, BlockType.NOTE) and b.text
            ]
            if not body_texts or len(slide.placeholders) < 2:
                continue

            body_placeholder = slide.placeholders[1]
            tf = body_placeholder.text_frame
            tf.text = body_texts[0]
            for extra in body_texts[1:]:
                p = tf.add_paragraph()
                p.text = extra
                p.level = 0

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
