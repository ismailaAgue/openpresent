"""
PPTX Export adapter — Technical Blueprint Section 3.5 / 7.

Revision (ADR-024): applies real typography and visual polish for the
first time. Previously, only the title's text COLOR was ever set —
font family, body text size, and every custom-layout textbox's font
were left completely unstyled (font.name=None, font.size=None),
meaning every deck rendered in PowerPoint's bare default look
regardless of the Theme chosen by the Design Engine. `Theme.font_set_id`
had existed since Phase 1 and was never actually read anywhere. This
was the real cause of "mediocre quality" feedback — a structure and
layout problem had already been solved (ADR-020/022/023), but the
underlying visual design was never built, only planned for.

Fixed here: a real font family per theme, consistent sizing across
every text element (title, body bullets, and every custom-layout
textbox), and a simple accent-colored bar under each title as a
visual anchor — the smallest addition that makes a deck read as
"designed" rather than "default," without adding image/AI complexity.
"""

import io
import re
from backend.ports.export import ExportPort
from backend.models.recipe import Recipe, BlockType
from backend.layout.layout_classifier import PROCESS_BULLET_PATTERN

# ADR-059 — deliberately broader than layout_classifier's own
# STATISTIC_PATTERN, and used only here, not shared with it. The
# classifier's pattern is tuned to avoid FALSE POSITIVES (only $ and %
# count, so an incidental "3 things" in an ordinary bullet doesn't
# wrongly flip a slide into the statistics layout — see that module's
# own comments). But once a slide IS already a statistics slide, chip
# rendering needs to find whatever number is actually there, including
# a bare magnitude like "415K" with no $ or % attached — a real gap
# found by rendering a real chip and seeing "415K happy customers
# served" overflow off the card because nothing matched at all and the
# whole string fell back to one oversized line. No trailing \b: it was
# tried first, but "%"/"K"/"M"/"B" followed by whitespace is a
# non-word-to-non-word transition (no boundary exists there), which
# forced the engine to backtrack the optional suffix OUT of the match
# — "90% client..." matched only "90", stranding the "%" in the label
# half. Also confirmed by rendering, not just reasoned about.
#
# ADR-062 — extended with an optional leading +/- sign: a climate-
# style stat like "+1.1C global warming since 1880" left the "+"
# stranded in the label half ("+C GLOBAL WARMING..."), separated from
# the "1.1" it actually belongs to — caught by rendering a real
# editorial-theme stats slide, not assumed. A trailing unit letter
# (the "C" in "1.1C", "°C", "x", etc.) is a stated, deliberate
# non-fix: consuming arbitrary trailing letters risks eating the start
# of the label text itself ("10 years" -> "10y" + "ears"), a worse
# failure mode than a unit letter staying in the label.
CHIP_NUMBER_PATTERN = re.compile(r"[+-]?\$?\d[\d,]*(\.\d+)?\s*[%KMB]?")


def _tint(rgb: tuple[int, int, int], amount: float) -> tuple[int, int, int]:
    """Lightens a color by blending it toward white. amount=0 returns
    the original color, amount=1 returns white. Used for stat-chip
    background boxes (ADR-059) — the chip needs to be visibly tinted
    with the theme's accent color without being so saturated that the
    dark number/label text inside it becomes hard to read."""
    r, g, b = rgb
    return (
        int(r + (255 - r) * amount),
        int(g + (255 - g) * amount),
        int(b + (255 - b) * amount),
    )

_COLOR_SETS = {
    "neutral": {
        "title": (0x22, 0x22, 0x22), "accent": (0x2E, 0x5C, 0x8A),
        "background": (0xF7, 0xF7, 0xF4), "text": (0x33, 0x33, 0x33),
        "corner_style": "circle", "stat_chip": False,
    },
    "blue_academic": {
        "title": (0x1B, 0x3A, 0x5C), "accent": (0xC8, 0x6B, 0x2E),
        "background": (0xF1, 0xF4, 0xF8), "text": (0x2A, 0x2A, 0x2A),
        "corner_style": "circle", "stat_chip": False,
    },
    # ADR-029 (presentation variety, spec Section 10)
    "warm_editorial": {
        "title": (0x3D, 0x25, 0x1A), "accent": (0xD9, 0x6C, 0x2E),
        "background": (0xFB, 0xF3, 0xEA), "text": (0x3D, 0x2E, 0x24),
        "corner_style": "circle", "stat_chip": False,
    },
    "modern_dark": {
        "title": (0xF2, 0xF2, 0xF0), "accent": (0x5D, 0xC9, 0xB0),
        "background": (0x1E, 0x21, 0x24), "text": (0xD8, 0xDA, 0xDC),
        "corner_style": "circle", "stat_chip": False,
    },
    # ADR-059 — four template themes, each modeled on a reference image
    # the person supplied. "corner_style"/"stat_chip"/"gradient_stops"
    # are new, theme-level visual controls (not just color swaps) that
    # _add_corner_decoration and _render_statistics_slide now read —
    # see those methods for exactly what each value changes.
    "gradient_violet": {
        # Reference: purple/blue/pink AI-themed deck — a soft gradient
        # blob in the corner and each key number in a colored "chip"
        # box, not plain text.
        "title": (0x24, 0x1B, 0x3D), "accent": (0xB0, 0x3D, 0xE8),
        "background": (0xFC, 0xFB, 0xFF), "text": (0x33, 0x2B, 0x44),
        "corner_style": "blob", "stat_chip": True,
        "gradient_stops": ((0x6A, 0x3D, 0xE8), (0xE8, 0x4D, 0xB8)),
    },
    "minimal_mono": {
        # Reference: grayscale/near-black minimal deck — no corner
        # decoration at all, no stat chips, everything rests on
        # typography and negative space alone.
        "title": (0x14, 0x14, 0x14), "accent": (0x6B, 0x6B, 0x6B),
        "background": (0xFF, 0xFF, 0xFF), "text": (0x2E, 0x2E, 0x2E),
        "corner_style": "none", "stat_chip": False,
    },
    "bold_violet_stats": {
        # Reference: bold black headlines on a violet/lavender palette
        # with punchy stat numbers — a plain accent circle (not a
        # gradient blob) and stat numbers in a strong accent color,
        # but no chip box around them; the boldness comes from
        # typography weight and color, not extra shapes.
        "title": (0x18, 0x14, 0x24), "accent": (0x6E, 0x4A, 0xE8),
        "background": (0xF7, 0xF5, 0xFD), "text": (0x2A, 0x24, 0x38),
        "corner_style": "circle", "stat_chip": False,
    },
    "clean_saas_blue": {
        # Reference: OpenPresent's own concept mockup — light, airy,
        # card-like, restrained blue accent. Closest in spirit to
        # "neutral" but with a distinct blue identity and no corner
        # decoration, matching that mockup's clean edges.
        "title": (0x1A, 0x22, 0x33), "accent": (0x3B, 0x6E, 0xF6),
        "background": (0xFA, 0xFB, 0xFD), "text": (0x33, 0x3B, 0x47),
        "corner_style": "none", "stat_chip": False,
    },
    # ADR-062 — an "editorial" theme, built in direct response to a
    # side-by-side critique against a competitor deck: full-bleed
    # cropped photography, a serif display headline, kicker/footer
    # metadata on every slide, and statistics rendered as a stacked
    # sidebar panel rather than a row of chips. This is not a color
    # palette variation of the existing render logic — it dispatches
    # to a genuinely different set of renderers
    # (_render_editorial_title_slide/_render_editorial_content_slide/
    # _render_editorial_stats_slide), selected via layout_style below.
    "editorial_cream": {
        "title": (0x16, 0x34, 0x2A), "accent": (0xC1, 0x5F, 0x2F),
        "background": (0xF6, 0xF2, 0xE9), "text": (0x2B, 0x2B, 0x28),
        "corner_style": "none", "stat_chip": False,
        "layout_style": "editorial",
    },
}

# Office-native fonts only — guarantees consistent rendering across
# any real PowerPoint install, rather than a font that might not be
# present and silently falls back to something else.
_FONT_SETS = {
    "sans": "Calibri",
    "serif": "Cambria",
}

TITLE_SLIDE_FONT_SIZE = 54
CONTENT_TITLE_FONT_SIZE = 40
BODY_FONT_SIZE = 20

# "Measure the title before deciding the layout" — a long title in a
# narrower column (when sharing the slide with an image) was measured
# overflowing past the bottom of the slide at a fixed 54pt. Scaling
# font size down by title length is simple, deterministic, and doesn't
# depend on whether a given renderer honors PowerPoint's autofit
# (LibreOffice's autofit fidelity is inconsistent; a computed size is
# correct everywhere, always).
def _fitting_title_font_size(title: str, base_size: int, narrow_column: bool) -> int:
    length = len(title)
    if narrow_column:
        if length > 70:
            ratio = 0.44
        elif length > 45:
            ratio = 0.56
        elif length > 28:
            ratio = 0.67
        else:
            ratio = 1.0
    else:
        if length > 70:
            ratio = 0.6
        elif length > 45:
            ratio = 0.8
        else:
            ratio = 1.0
    return max(18, int(base_size * ratio))  # never shrink below a readable floor


class PptxExportAdapter(ExportPort):
    def format_id(self) -> str:
        return "pptx"

    def export(self, recipe: Recipe) -> bytes:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.enum.shapes import MSO_SHAPE
        except ImportError as e:
            raise RuntimeError(
                "python-pptx is required for PPTX export. Install with: "
                "pip install python-pptx --break-system-packages"
            ) from e

        colors = _COLOR_SETS.get(recipe.theme.color_set_id, _COLOR_SETS["neutral"])
        title_color = RGBColor(*colors["title"])
        accent_color = RGBColor(*colors["accent"])
        background_color = RGBColor(*colors["background"])
        text_color = RGBColor(*colors.get("text", colors["title"]))
        font_name = _FONT_SETS.get(recipe.theme.font_set_id, _FONT_SETS["sans"])
        corner_style = colors.get("corner_style", "circle")
        stat_chip = colors.get("stat_chip", False)
        gradient_stops = colors.get("gradient_stops")

        # Lazy import to avoid a circular dependency — registry.py
        # imports this module to construct PptxExportAdapter, so this
        # module can't import registry at the top level.
        from backend.adapters import registry as _registry
        media = _registry.get_media_adapter()

        ctx = _RenderContext(
            Inches=Inches, Pt=Pt, Emu=Emu, RGBColor=RGBColor, PP_ALIGN=PP_ALIGN,
            MSO_SHAPE=MSO_SHAPE, title_color=title_color, accent_color=accent_color,
            background_color=background_color, text_color=text_color, font_name=font_name, media=media,
            corner_style=corner_style, stat_chip=stat_chip, gradient_stops=gradient_stops,
            accent_rgb=colors["accent"], layout_style=colors.get("layout_style", "default"),
        )

        prs = Presentation()
        content_layout = prs.slide_layouts[1]
        # "Title Only" layout (index 5 in the default template) — gives
        # a title placeholder with an empty body, so every specialized
        # layout (title slide, statistics, comparison, process, and
        # now bullet+image) can place shapes at explicitly-computed,
        # verified-non-overlapping positions instead of relying on a
        # fixed template placeholder that doesn't account for an image
        # sharing the slide (ADR-026 — see module docstring).
        title_only_layout = prs.slide_layouts[5]

        for i, slide_data in enumerate(sorted(recipe.outline.slides, key=lambda s: s.order)):
            is_title_slide = (i == 0)

            if ctx.layout_style == "editorial":
                # ADR-062 — a genuinely different renderer set, not a
                # color variation of the default one. Comparison/
                # process layout_types fall back to the same content
                # renderer everything else here uses (a stated scope
                # limit — see this ADR's own notes on what's covered).
                total_slides = len(recipe.outline.slides)
                deck_title = recipe.outline.slides[0].title if recipe.outline.slides else ""
                if is_title_slide:
                    slide = self._render_editorial_title_slide(prs, slide_data, ctx)
                else:
                    body_texts = [b.text for b in slide_data.content_blocks if b.type == BlockType.BULLET and b.text]
                    if slide_data.layout_type == "statistics" and body_texts:
                        slide = self._render_editorial_stats_slide(
                            prs, slide_data, body_texts, ctx, section_number=i, deck_title=deck_title, index=i + 1, total=total_slides
                        )
                    else:
                        slide = self._render_editorial_content_slide(
                            prs, slide_data, body_texts, ctx, section_number=i, deck_title=deck_title, index=i + 1, total=total_slides
                        )
                notes_text = "\n".join(b.text for b in slide_data.content_blocks if b.type == BlockType.NOTE and b.text)
                if notes_text and slide is not None:
                    slide.notes_slide.notes_text_frame.text = notes_text
                continue

            if is_title_slide:
                slide = self._render_title_slide(prs, title_only_layout, slide_data, ctx)
            else:
                # FIX (ADR-029): previously included BlockType.NOTE text
                # in body_texts, meaning speaker notes were rendered as
                # visible bullets on every slide instead of going into
                # PowerPoint's actual notes pane. Bullets only, now —
                # notes are set on slide.notes_slide below instead.
                body_texts = [
                    b.text for b in slide_data.content_blocks
                    if b.type == BlockType.BULLET and b.text
                ]

                if slide_data.layout_type == "comparison" and body_texts:
                    slide = self._render_comparison_slide(prs, title_only_layout, slide_data.title, body_texts, ctx)
                elif slide_data.layout_type == "process" and body_texts:
                    slide = self._render_process_slide(prs, title_only_layout, slide_data.title, body_texts, ctx)
                elif slide_data.layout_type == "statistics" and body_texts:
                    slide = self._render_statistics_slide(prs, title_only_layout, slide_data.title, body_texts, ctx)
                else:
                    image_query = getattr(slide_data, "image_query", None)
                    slide = self._render_bullet_slide(prs, content_layout, title_only_layout, slide_data.title,
                                                       body_texts, ctx, media=ctx.media, image_query=image_query)

            notes_text = "\n".join(
                b.text for b in slide_data.content_blocks if b.type == BlockType.NOTE and b.text
            )
            if notes_text and slide is not None:
                slide.notes_slide.notes_text_frame.text = notes_text

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # -- shared helpers ---------------------------------------------------

    def _add_title_with_accent(self, slide, ctx, title):
        """Sets a styled content-slide title and adds a thin accent-
        colored bar beneath it. (The pptx design skill flags accent
        bars as a common AI-slop pattern — noted, but kept here per
        explicit product decision: it's a deliberate visual identity
        choice, not an oversight.) Used by every layout so all of them
        read as one consistent product."""
        slide.shapes.title.text = title
        fitted_size = _fitting_title_font_size(title, CONTENT_TITLE_FONT_SIZE, narrow_column=False)
        ctx.style_run(slide.shapes.title.text_frame.paragraphs[0],
                       size=fitted_size, color=ctx.title_color, bold=True)

        title_shape = slide.shapes.title
        bar_top = title_shape.top + title_shape.height + ctx.Emu(45000)
        bar = slide.shapes.add_shape(
            ctx.MSO_SHAPE.RECTANGLE, title_shape.left, bar_top, ctx.Inches(1.4), ctx.Pt(4)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ctx.accent_color
        bar.line.fill.background()
        bar.shadow.inherit = False

    def _apply_background(self, slide, ctx):
        """Tier 1 visual improvement: a real background fill instead of
        plain white on every slide, per the theme's background color.
        Cheap, pure-rules, and directly addresses one of the clearest
        gaps versus the reference decks (which never use plain white)."""
        if ctx.background_color is None:
            return
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = ctx.background_color

    def _add_corner_decoration(self, slide, ctx, prs, small=False):
        """Tier 1 visual improvement: a simple geometric accent shape
        in a slide corner — cheap to draw, adds visual interest without
        needing any image, and is a real (if modest) step toward the
        "don't create text-only slides" guidance, ahead of Tier 2's
        real image integration.

        ADR-059 — three variants now, chosen per-theme rather than one
        fixed look for every theme: "none" skips this entirely (for
        minimal/restrained themes where an extra shape would work
        against the point), "blob" draws a larger soft-edged gradient
        shape (for themes modeled on gradient-heavy reference decks),
        and "circle" is the original plain accent-colored oval,
        unchanged, still the default for every pre-existing theme.

        ADR-061 — small=True gives every non-title slide a much
        smaller version of the same shape, tucked into the bottom-
        right corner (title slides keep the large top-left version,
        unchanged), rather than skipping decoration on every slide
        after the first. This was a real, reported gap: the reference
        decks carry a consistent visual identity mark across every
        slide, not just the cover, and a deck that's decorated only on
        slide 1 reads as "handmade"/inconsistent by slide 3.

        Takes `prs` explicitly (not assumed to be 13.333x7.5in
        widescreen) — python-pptx's default blank Presentation() is
        actually 10x7.5in (4:3); the first version of this method
        hardcoded 16:9 dimensions for the small-variant's bottom-right
        positioning, which would have placed it off-canvas entirely.
        Caught before shipping, not after."""
        if ctx.corner_style == "none":
            return
        slide_width, slide_height = prs.slide_width, prs.slide_height
        if ctx.corner_style == "blob" and ctx.gradient_stops:
            size = ctx.Inches(1.0) if small else ctx.Inches(2.6)
            offset = ctx.Inches(-0.35) if small else ctx.Inches(-0.8)
            left = slide_width - size - offset if small else offset
            top = slide_height - size - offset if small else offset
            shape = slide.shapes.add_shape(ctx.MSO_SHAPE.OVAL, left, top, size, size)
            shape.fill.gradient()
            shape.fill.gradient_angle = 45.0  # must be set before stop colors — see module notes
            stops = shape.fill.gradient_stops
            stops[0].color.rgb = ctx.RGBColor(*ctx.gradient_stops[0])
            stops[0].position = 0.0
            stops[1].color.rgb = ctx.RGBColor(*ctx.gradient_stops[1])
            stops[1].position = 1.0
            shape.line.fill.background()
            shape.shadow.inherit = False
            return

        size = ctx.Inches(0.35) if small else ctx.Inches(0.9)
        offset = ctx.Inches(-0.12) if small else ctx.Inches(-0.3)
        left = slide_width - size - offset if small else offset
        top = slide_height - size - offset if small else offset
        shape = slide.shapes.add_shape(ctx.MSO_SHAPE.OVAL, left, top, size, size)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ctx.accent_color
        shape.line.fill.background()
        shape.shadow.inherit = False

    def _disable_native_bullet(self, paragraph):
        """ADR-061 — python-pptx 1.0.2 has no paragraph.bullet API (any
        version's docs referencing it are for a different/newer
        release — confirmed against the actual installed version, not
        assumed), so the only way to stop a placeholder's inherited
        bullet character from showing is direct XML: insert <a:buNone/>
        into the paragraph's <a:pPr>. Without this, _add_colored_bullets'
        own accent-colored marker run would show up ALONGSIDE the
        native bullet instead of replacing it — a real, worse-than-
        original double-bullet bug this was caught by testing, not
        assumed to work from the (wrong) high-level API guess."""
        from pptx.oxml.ns import qn
        pPr = paragraph._p.get_or_add_pPr()
        for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
            existing = pPr.find(qn(tag))
            if existing is not None:
                pPr.remove(existing)
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))

    def _add_colored_bullets(self, text_frame, texts, ctx, size=BODY_FONT_SIZE):
        """ADR-061 — replaces the default black round-dot bullet (or,
        on the native-placeholder path, whatever the inherited layout
        happened to use) with an accent-colored square marker as its
        own text run, ahead of the bullet text itself — a small,
        cheap, but real "designed, not default" signal the reference
        decks all share (every one of them uses a colored marker, none
        use a plain black dot)."""
        if not texts:
            return
        text_frame.word_wrap = True
        for i, txt in enumerate(texts):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = ""
            self._disable_native_bullet(p)
            marker_run = p.add_run()
            marker_run.text = "■  "
            ctx.style_run_direct(marker_run, size=size, color=ctx.accent_color, bold=False)
            text_run = p.add_run()
            text_run.text = txt
            ctx.style_run_direct(text_run, size=size, color=ctx.text_color, bold=False)
            p.space_after = ctx.Pt(10)

    def _add_picture_cover_crop(self, slide, ctx, image_bytes, left, top, width, height):
        """True edge-to-edge 'full bleed' image fill, preserving aspect
        ratio via cropping (never stretching/distorting) — the single
        biggest, most concrete visual gap identified against the
        reference decks (ADR-062): every image in the deck up to this
        point was placed inside a bounded, padded rectangle with
        whitespace around it, never filling a region edge-to-edge the
        way an editorial layout's photography does. Computes real crop
        fractions from the image's actual pixel dimensions (via PIL,
        already a project dependency — reportlab pulls it in) against
        the target box's aspect ratio, rather than letting
        add_picture's default stretch-to-fit distort the photo."""
        import io as _io
        from PIL import Image as _Image
        try:
            with _Image.open(_io.BytesIO(image_bytes)) as im:
                img_w, img_h = im.size
        except Exception:
            img_w, img_h = 4, 3  # fallback aspect if the bytes aren't readable as an image
        box_ratio = width / height
        img_ratio = img_w / img_h
        pic = slide.shapes.add_picture(_io.BytesIO(image_bytes), left, top, width, height)
        if img_ratio > box_ratio:
            # image is relatively wider than the box — crop left/right
            visible_frac = box_ratio / img_ratio
            crop = (1 - visible_frac) / 2
            pic.crop_left = crop
            pic.crop_right = crop
        elif img_ratio < box_ratio:
            # image is relatively taller than the box — crop top/bottom
            visible_frac = img_ratio / box_ratio
            crop = (1 - visible_frac) / 2
            pic.crop_top = crop
            pic.crop_bottom = crop
        pic.line.fill.background()
        return pic

    def _add_kicker(self, slide, ctx, left, top, width, section_number, label):
        """Small-caps 'SECTION N — LABEL' line with a thin trailing
        rule, matching the reference decks' recurring editorial
        identity element (ADR-062) — every content slide gets one, not
        just the cover, the same "consistent across every slide, not
        just the first" principle ADR-061's corner decoration and
        colored bullets already established for the default themes.

        Truncates on a WORD boundary, not a character count — the
        first version sliced the raw title to a fixed character
        length, which cut mid-word ("...RUNNING A F" instead of
        "...RUNNING A FEVER") on anything longer than the limit.
        Caught by rendering, not by re-reading the slice logic."""
        words = label.split()
        short = []
        length = 0
        for w in words:
            if length + len(w) + 1 > 28:
                break
            short.append(w)
            length += len(w) + 1
        text = f"{section_number:02d} — {' '.join(short).upper()}"
        box = slide.shapes.add_textbox(left, top, width, ctx.Inches(0.3))
        p = box.text_frame.paragraphs[0]
        p.text = text
        ctx.style_run(p, size=10, color=ctx.accent_color, bold=True)
        # Slightly wider letter-spacing isn't exposed by python-pptx's
        # high-level API; the small size + bold + accent color alone
        # is enough to read as a "kicker" label rather than body text.
        return box

    def _add_footer(self, slide, ctx, prs, deck_title, index, total, skip_left=False):
        """Deck-title / page-counter row at the bottom of every slide
        — the other half of the recurring editorial identity, alongside
        the kicker. Both numbers are computed purely from the slide's
        own position in the deck, no new content field required.

        skip_left=True omits the deck-title label specifically — used
        on content slides with a full-bleed image on the left half,
        where the deck-title text would otherwise land directly on top
        of that image's own attribution caption. Caught by rendering a
        real image-bearing slide and seeing the two overlap, not
        assumed from the layout math."""
        slide_width, slide_height = prs.slide_width, prs.slide_height
        margin = ctx.Inches(0.6)
        footer_top = slide_height - ctx.Inches(0.5)
        if not skip_left:
            left_box = slide.shapes.add_textbox(margin, footer_top, ctx.Inches(3.5), ctx.Inches(0.3))
            lp = left_box.text_frame.paragraphs[0]
            lp.text = deck_title.upper()[:40]
            ctx.style_run(lp, size=8, color=ctx.text_color, bold=False)

        right_box = slide.shapes.add_textbox(slide_width - margin - ctx.Inches(1.5), footer_top, ctx.Inches(1.5), ctx.Inches(0.3))
        rp = right_box.text_frame.paragraphs[0]
        rp.text = f"{index:02d} / {total:02d}"
        rp.alignment = ctx.PP_ALIGN.RIGHT
        ctx.style_run(rp, size=8, color=ctx.text_color, bold=False)

    def _render_editorial_title_slide(self, prs, slide_data, ctx):
        """Cover slide for editorial-layout themes (ADR-062): a
        full-bleed image filling exactly one half of the canvas edge
        to edge (via _add_picture_cover_crop, never stretched), a
        kicker line, a large serif title, and an optional subtitle
        paragraph beneath it — deliberately NOT the boxed/padded/
        captioned image treatment every other theme uses, matching
        what the reference decks' cover slides actually do."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # fully blank — no inherited placeholders to fight
        self._apply_background(slide, ctx)
        slide_width, slide_height = prs.slide_width, prs.slide_height
        margin = ctx.Inches(0.55)

        image_result = self._maybe_fetch_image(ctx.media, getattr(slide_data, "image_query", None), ctx)
        text_width = slide_width // 2 if image_result else slide_width - (2 * margin)

        if image_result:
            self._add_picture_cover_crop(slide, ctx, image_result.image_bytes, slide_width // 2, 0, slide_width // 2, slide_height)

        kicker_top = ctx.Inches(0.5)
        self._add_kicker(slide, ctx, margin, kicker_top, text_width - margin, 1, "A PRIMER")

        title_top = ctx.Inches(1.3)
        title_box = slide.shapes.add_textbox(margin, title_top, text_width - margin, ctx.Inches(2.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        fitted_size = _fitting_title_font_size(slide_data.title, 48, narrow_column=bool(image_result))
        tf.text = slide_data.title
        ctx.style_run(tf.paragraphs[0], size=fitted_size, color=ctx.title_color, bold=False)  # serif display weight, not bold

        return slide

    def _render_editorial_content_slide(self, prs, slide_data, body_texts, ctx, section_number, deck_title, index, total):
        """Regular content slide for editorial-layout themes (ADR-062):
        kicker + serif title + real prose paragraphs (falling back to
        minimal dash-marker bullets only for genuinely fragment-like
        content — see _looks_like_prose), plus a full-bleed image
        filling the OTHER half when one's available, alternating which
        side has the text vs the image isn't attempted (always image-
        left/text-right, a stated simplification — see this method's
        scope note in the module's ADR-062 entry)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._apply_background(slide, ctx)
        slide_width, slide_height = prs.slide_width, prs.slide_height
        margin = ctx.Inches(0.6)

        image_result = self._maybe_fetch_image(ctx.media, getattr(slide_data, "image_query", None), ctx)
        if image_result:
            self._add_picture_cover_crop(slide, ctx, image_result.image_bytes, 0, 0, slide_width // 2, slide_height)
            if image_result.attribution:
                cap_box = slide.shapes.add_textbox(
                    ctx.Inches(0.3), slide_height - ctx.Inches(0.55), slide_width // 2 - ctx.Inches(0.6), ctx.Inches(0.4)
                )
                cp = cap_box.text_frame.paragraphs[0]
                cp.text = image_result.attribution
                ctx.style_run(cp, size=9, color=ctx.RGBColor(0xFF, 0xFF, 0xFF), bold=False)
            text_left = slide_width // 2 + margin
            text_width = slide_width // 2 - margin - ctx.Inches(0.3)
        else:
            text_left = margin
            text_width = slide_width - (2 * margin)

        self._add_kicker(slide, ctx, text_left, ctx.Inches(0.45), text_width, section_number, slide_data.title)

        title_box = slide.shapes.add_textbox(text_left, ctx.Inches(0.85), text_width, ctx.Inches(1.4))
        tf = title_box.text_frame
        tf.word_wrap = True
        fitted_size = _fitting_title_font_size(slide_data.title, 30, narrow_column=True)
        tf.text = slide_data.title
        ctx.style_run(tf.paragraphs[0], size=fitted_size, color=ctx.title_color, bold=False)

        body_top = ctx.Inches(2.3)
        body_height = slide_height - body_top - ctx.Inches(0.9)
        body_box = slide.shapes.add_textbox(text_left, body_top, text_width, body_height)
        btf = body_box.text_frame
        btf.word_wrap = True
        if self._looks_like_prose(body_texts):
            p = btf.paragraphs[0]
            p.text = body_texts[0]
            ctx.style_run(p, size=13, color=ctx.text_color)
            p.space_after = ctx.Pt(12)
            for extra in body_texts[1:]:
                np = btf.add_paragraph()
                np.text = extra
                ctx.style_run(np, size=13, color=ctx.text_color)
                np.space_after = ctx.Pt(12)
        else:
            for i, txt in enumerate(body_texts):
                p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                p.text = ""
                self._disable_native_bullet(p)
                marker = p.add_run()
                marker.text = "—  "
                ctx.style_run_direct(marker, size=13, color=ctx.accent_color)
                text_run = p.add_run()
                text_run.text = txt
                ctx.style_run_direct(text_run, size=13, color=ctx.text_color)
                p.space_after = ctx.Pt(10)

        self._add_footer(slide, ctx, prs, deck_title, index, total, skip_left=bool(image_result))
        return slide

    @staticmethod
    def _looks_like_prose(texts: list[str]) -> bool:
        """Same rule DocumentDocxExportAdapter/DocumentPdfExportAdapter
        already use to distinguish real sentences from bullet
        fragments (ADR-054) — reused here rather than reinvented, since
        it's the same underlying question: does this content read as
        connected prose or as a list of fragments?"""
        return len(texts) >= 1 and all(t.rstrip().endswith((".", "!", "?")) for t in texts)

    def _render_editorial_stats_slide(self, prs, slide_data, body_texts, ctx, section_number, deck_title, index, total):
        """Statistics as a vertical stacked sidebar panel — a real
        design element, not text arranged in a row (ADR-062). This is
        the single most-cited gap in the reference-deck critique:
        'OpenPresent treats statistics as text; the reference deck
        treats them as design elements.' Each stat's number is pulled
        out and set large or with an accent color; the label sits
        small and muted beneath it; a thin rule separates each row."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._apply_background(slide, ctx)
        slide_width, slide_height = prs.slide_width, prs.slide_height
        margin = ctx.Inches(0.6)

        panel_left = int(slide_width * 0.58)
        panel_width = slide_width - panel_left - margin
        panel_top = ctx.Inches(0.55)
        panel_height = slide_height - panel_top - ctx.Inches(0.9)

        panel = slide.shapes.add_shape(ctx.MSO_SHAPE.RECTANGLE, panel_left, panel_top, panel_width, panel_height)
        panel.fill.solid()
        panel.fill.fore_color.rgb = ctx.RGBColor(*_tint(ctx.accent_rgb, 0.92)) if ctx.accent_rgb else ctx.background_color
        panel.line.fill.background()
        panel.shadow.inherit = False
        top_bar = slide.shapes.add_shape(ctx.MSO_SHAPE.RECTANGLE, panel_left, panel_top, panel_width, ctx.Inches(0.05))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = ctx.title_color
        top_bar.line.fill.background()
        top_bar.shadow.inherit = False

        stats = body_texts[:4]
        row_height = panel_height // max(1, len(stats))
        pad = ctx.Inches(0.25)
        for i, stat_text in enumerate(stats):
            row_top = panel_top + (i * row_height)
            match = CHIP_NUMBER_PATTERN.search(stat_text)
            if match:
                number_part = match.group(0).strip()
                label_part = (stat_text[:match.start()] + stat_text[match.end():]).strip(" -:,.")
            else:
                number_part, label_part = stat_text, ""
            num_box = slide.shapes.add_textbox(panel_left + pad, row_top + ctx.Inches(0.15), panel_width - (2 * pad), ctx.Inches(0.55))
            np_ = num_box.text_frame.paragraphs[0]
            np_.text = number_part
            number_size = 26 if len(number_part) <= 8 else 18
            ctx.style_run(np_, size=number_size, color=(ctx.accent_color if i % 2 == 0 else ctx.title_color), bold=True)
            if label_part:
                lbl_box = slide.shapes.add_textbox(panel_left + pad, row_top + ctx.Inches(0.68), panel_width - (2 * pad), ctx.Inches(0.4))
                lp = lbl_box.text_frame.paragraphs[0]
                lp.text = label_part.upper()
                ctx.style_run(lp, size=9, color=ctx.text_color)
            if i > 0:
                divider = slide.shapes.add_shape(ctx.MSO_SHAPE.RECTANGLE, panel_left + pad, row_top, panel_width - (2 * pad), ctx.Emu(9525))
                divider.fill.solid()
                divider.fill.fore_color.rgb = ctx.RGBColor(*_tint(ctx.accent_rgb, 0.6)) if ctx.accent_rgb else ctx.text_color
                divider.line.fill.background()
                divider.shadow.inherit = False

        text_left = margin
        text_width = panel_left - margin - ctx.Inches(0.4)
        self._add_kicker(slide, ctx, text_left, ctx.Inches(0.45), text_width, section_number, slide_data.title)
        title_box = slide.shapes.add_textbox(text_left, ctx.Inches(0.85), text_width, ctx.Inches(1.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        fitted_size = _fitting_title_font_size(slide_data.title, 32, narrow_column=True)
        tf.text = slide_data.title
        ctx.style_run(tf.paragraphs[0], size=fitted_size, color=ctx.title_color, bold=False)

        self._add_footer(slide, ctx, prs, deck_title, index, total)
        return slide

    # -- layout renderers ------------------------------------------------


    def _render_bullet_slide(self, prs, content_layout, title_only_layout, title, body_texts, ctx, media=None, image_query=None):
        image_result = self._maybe_fetch_image(media, image_query, ctx)

        if image_result and body_texts:
            # Uses title_only_layout + manually-positioned shapes, same
            # as every other specialized layout — NOT the native content
            # placeholder. Resizing an inherited placeholder's left/width
            # without also setting top/height corrupts its position
            # (found via real testing — ADR-026): the body placeholder
            # serialized with height=0 and top=0, landing directly under
            # the title instead of below it. Manual textboxes with fully
            # explicit geometry sidestep that class of bug entirely.
            return self._render_bullet_slide_with_image(prs, title_only_layout, title, body_texts, image_result, ctx)

        # No image (or no media adapter configured) — the plain,
        # unchanged path using the native placeholder, which is never
        # resized and so was never affected by the bug above.
        slide = prs.slides.add_slide(content_layout)
        self._apply_background(slide, ctx)
        self._add_corner_decoration(slide, ctx, prs, small=True)
        self._add_title_with_accent(slide, ctx, title)

        if not body_texts or len(slide.placeholders) < 2:
            return slide

        body_placeholder = slide.placeholders[1]
        self._add_colored_bullets(body_placeholder.text_frame, body_texts, ctx)
        return slide

    def _render_bullet_slide_with_image(self, prs, title_only_layout, title, body_texts, image_result, ctx):
        """Text on the left, image on the right — both positioned
        explicitly below the title's REAL bottom edge (read after the
        title is actually placed, not assumed), guaranteeing no
        title/content/image overlap by construction rather than by
        coincidence (ADR-026 — this is the fix for the confirmed P0
        collision bug found via real generated output)."""
        slide = prs.slides.add_slide(title_only_layout)
        self._apply_background(slide, ctx)
        self._add_corner_decoration(slide, ctx, prs, small=True)
        self._add_title_with_accent(slide, ctx, title)

        title_shape = slide.shapes.title
        # Real measured bottom edge of the title, plus clearance for
        # the accent bar beneath it — not a guessed constant.
        content_top = title_shape.top + title_shape.height + ctx.Inches(0.3)

        slide_width, slide_height = prs.slide_width, prs.slide_height
        margin = ctx.Inches(0.5)
        gutter = ctx.Inches(0.3)
        bottom_margin = ctx.Inches(0.4)
        image_width = ctx.Inches(4.0)
        text_width = slide_width - margin - gutter - image_width - margin
        content_height = slide_height - content_top - bottom_margin

        text_box = slide.shapes.add_textbox(margin, content_top, text_width, content_height)
        self._add_colored_bullets(text_box.text_frame, body_texts, ctx)

        image_left = margin + text_width + gutter
        pic_bottom = self._add_picture_capped(slide, ctx, image_result.image_bytes,
                                               image_left, content_top, image_width, content_height)
        self._add_attribution_caption(slide, ctx, image_result, image_left, pic_bottom, image_width)
        return slide

    def _add_picture_capped(self, slide, ctx, image_bytes, left, top, max_width, max_height):
        """Adds a picture sized to max_width, then shrinks it
        proportionally (preserving aspect ratio) if its natural height
        would exceed max_height — guarantees an image can never extend
        past its allotted region regardless of its native aspect ratio,
        which is what actually prevents overlap, not just a width cap.
        Returns the picture's bottom edge (EMU) so callers can place an
        attribution caption directly beneath it without recomputing."""
        import io as _io
        pic = slide.shapes.add_picture(_io.BytesIO(image_bytes), left, top, width=max_width)
        if pic.height > max_height:
            scale = max_height / pic.height
            pic.height = int(pic.height * scale)
            pic.width = int(pic.width * scale)
            pic.left = left  # re-anchor left edge; only height/width shrank
        return pic.top + pic.height

    def _add_attribution_caption(self, slide, ctx, image_result, left, top, width):
        """Tiny caption beneath an image — only rendered when the
        provider's license requires visible credit (Wikimedia Commons
        chiefly; Unsplash/Pexels/Pixabay set attribution=None since
        their free-use terms don't require an on-slide credit). Spec
        Section 8: 'attribution handling where required.'"""
        if not getattr(image_result, "attribution", None):
            return
        caption = slide.shapes.add_textbox(left, top, width, ctx.Inches(0.3))
        tf = caption.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = image_result.attribution[:120]
        p.alignment = ctx.PP_ALIGN.LEFT
        ctx.style_run(p, size=8, color=ctx.text_color)

    def _maybe_fetch_image(self, media, image_query, ctx=None):
        """Never raises, never blocks the deck on a failed/missing
        image — same graceful-degradation discipline as the AI Port
        (Constitution Principle 3, extended to media in ADR-025).
        Tracks used image_ids on ctx (ADR-029) so the same photo never
        appears twice in one deck — the router picks the next-best
        scored candidate instead."""
        if media is None or image_query is None:
            return None
        try:
            if not media.is_available():
                return None
            exclude = ctx.used_image_ids if ctx is not None else None
            result = media.search_image(image_query, exclude_ids=exclude)
            if result is not None and ctx is not None:
                ctx.used_image_ids.add(result.image_id)
            return result
        except Exception:
            return None

    def _render_title_slide(self, prs, title_only_layout, slide_data, ctx):
        """Title text and (optional) image are both positioned
        explicitly, in the same left-column/right-column scheme as
        every other layout — NOT the default 'Title Slide' template's
        fixed placeholder positions, which is what caused the
        confirmed P0 collision bug: that template's title box sits at
        a fixed vertical position regardless of content, and the old
        code placed the image at an independently-chosen fixed
        position with no check that the two didn't overlap. They did,
        on every single deck. Fixed here by computing both boxes from
        the same coordinate scheme, so overlap is structurally
        impossible rather than avoided by a lucky constant (ADR-026)."""
        slide = prs.slides.add_slide(title_only_layout)
        self._apply_background(slide, ctx)
        self._add_corner_decoration(slide, ctx, prs)

        image_result = self._maybe_fetch_image(ctx.media, getattr(slide_data, "image_query", None), ctx)
        image_bytes = image_result.image_bytes if image_result else None
        slide_width, slide_height = prs.slide_width, prs.slide_height
        margin = ctx.Inches(0.6)

        if image_bytes:
            gutter = ctx.Inches(0.4)
            image_width = ctx.Inches(4.2)
            title_width = slide_width - margin - gutter - image_width - margin
            title_top = int(slide_height * 0.28)
            title_height = int(slide_height * 0.62)  # generous — most of the remaining vertical space

            title_box = slide.shapes.add_textbox(margin, title_top, title_width, title_height)
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.text = slide_data.title
            fitted_size = _fitting_title_font_size(slide_data.title, TITLE_SLIDE_FONT_SIZE, narrow_column=True)
            ctx.style_run(tf.paragraphs[0], size=fitted_size, color=ctx.title_color, bold=True)

            image_left = margin + title_width + gutter
            image_top = ctx.Inches(0.6)
            image_max_height = slide_height - image_top - ctx.Inches(0.6)
            pic_bottom = self._add_picture_capped(slide, ctx, image_bytes, image_left, image_top, image_width, image_max_height)
            self._add_attribution_caption(slide, ctx, image_result, image_left, pic_bottom, image_width)
        else:
            title_top = int(slide_height * 0.32)
            title_height = int(slide_height * 0.55)
            title_box = slide.shapes.add_textbox(margin, title_top, slide_width - (2 * margin), title_height)
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.text = slide_data.title
            fitted_size = _fitting_title_font_size(slide_data.title, TITLE_SLIDE_FONT_SIZE, narrow_column=False)
            ctx.style_run(tf.paragraphs[0], size=fitted_size, color=ctx.title_color, bold=True)
        return slide

    def _render_statistics_slide(self, prs, title_only_layout, title, body_texts, ctx):
        """Large, centered callouts arranged in a row instead of a
        bulleted list — appropriate for slides that are mostly a
        handful of key numbers (Layout Classifier: ADR-022).

        ADR-059 — themes with stat_chip=True (modeled on a reference
        deck that put every key number inside a colored card, not
        plain text) get that treatment here: the number portion is
        split out from the rest of the bullet text (via the same
        STATISTIC_PATTERN the layout classifier itself uses to decide
        a slide even qualifies as a statistics slide) and rendered
        large inside a tinted rounded-rectangle card, with whatever
        text remains as a smaller label beneath it. Themes with
        stat_chip=False keep the original plain-text behavior,
        unchanged."""
        slide = prs.slides.add_slide(title_only_layout)
        self._apply_background(slide, ctx)
        self._add_corner_decoration(slide, ctx, prs, small=True)
        self._add_title_with_accent(slide, ctx, title)

        stats = body_texts[:4]  # beyond 4, a row stops being readable at a glance
        slide_width, slide_height = prs.slide_width, prs.slide_height
        margin = ctx.Inches(0.5)
        box_width = (slide_width - (2 * margin)) // len(stats)
        box_top = int(slide_height * 0.40)
        box_height = ctx.Inches(2.4)

        if not ctx.stat_chip:
            for idx, stat_text in enumerate(stats):
                left = margin + (idx * box_width)
                box = slide.shapes.add_textbox(left, box_top, box_width, box_height)
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = stat_text
                p.alignment = ctx.PP_ALIGN.CENTER
                ctx.style_run(p, size=22, color=ctx.accent_color, bold=True)
            return slide

        chip_gap = ctx.Inches(0.15)
        chip_fill = ctx.RGBColor(*_tint(ctx.accent_rgb, 0.85))  # a light tint, not the raw accent
        for idx, stat_text in enumerate(stats):
            left = margin + (idx * box_width) + chip_gap
            chip_width = box_width - (2 * chip_gap)

            chip = slide.shapes.add_shape(ctx.MSO_SHAPE.ROUNDED_RECTANGLE, left, box_top, chip_width, box_height)
            chip.fill.solid()
            chip.fill.fore_color.rgb = chip_fill
            chip.line.fill.background()
            chip.shadow.inherit = False

            match = CHIP_NUMBER_PATTERN.search(stat_text)
            if match:
                number_part = match.group(0).strip()
                label_part = (stat_text[:match.start()] + stat_text[match.end():]).strip(" -:,.")
            else:
                number_part, label_part = stat_text, ""

            pad = ctx.Inches(0.2)
            number_box = slide.shapes.add_textbox(left + pad, box_top + ctx.Inches(0.25), chip_width - (2 * pad), ctx.Inches(0.9))
            ntf = number_box.text_frame
            ntf.word_wrap = True
            np_ = ntf.paragraphs[0]
            np_.text = number_part
            np_.alignment = ctx.PP_ALIGN.CENTER
            # Long numbers ("$320,820M") need a smaller size than short
            # ones ("90%") to have any chance of fitting one line in a
            # card this width — same fitting-by-length idea as
            # _fitting_title_font_size, just a much smaller range.
            number_size = 26 if len(number_part) <= 6 else (20 if len(number_part) <= 10 else 16)
            ctx.style_run(np_, size=number_size, color=ctx.accent_color, bold=True)

            if label_part:
                label_box = slide.shapes.add_textbox(
                    left + pad, box_top + ctx.Inches(1.15), chip_width - (2 * pad), ctx.Inches(1.0)
                )
                lf = label_box.text_frame
                lf.word_wrap = True
                lp = lf.paragraphs[0]
                lp.text = label_part
                lp.alignment = ctx.PP_ALIGN.CENTER
                ctx.style_run(lp, size=12, color=ctx.text_color)
        return slide

    def _render_comparison_slide(self, prs, title_only_layout, title, body_texts, ctx):
        """Two side-by-side text columns instead of one bulleted list —
        for slides whose title signals a direct comparison ('X vs Y'),
        per the Layout Classifier (ADR-022).

        ADR-061 — each column now sits inside a light, tinted card
        (using the same _tint() helper the statistics chips use), not
        floating directly on the bare background — the reference decks
        never present two comparison columns as plain unbounded text;
        a visible card boundary is what actually reads as "a
        comparison," where two paragraphs with no visual separation
        just reads as two disconnected blocks of text."""
        slide = prs.slides.add_slide(title_only_layout)
        self._apply_background(slide, ctx)
        self._add_corner_decoration(slide, ctx, prs, small=True)
        self._add_title_with_accent(slide, ctx, title)

        midpoint = max(1, len(body_texts) // 2) if len(body_texts) > 1 else len(body_texts)
        left_items, right_items = body_texts[:midpoint], body_texts[midpoint:]

        slide_width, slide_height = prs.slide_width, prs.slide_height
        margin, gutter = ctx.Inches(0.5), ctx.Inches(0.3)
        column_width = (slide_width - (2 * margin) - gutter) // 2
        top = int(slide_height * 0.36)
        height = ctx.Inches(3.5)
        card_fill = ctx.RGBColor(*_tint(ctx.accent_rgb, 0.9)) if ctx.accent_rgb else None

        for items, left in ((left_items, margin), (right_items, margin + column_width + gutter)):
            if card_fill is not None:
                card = slide.shapes.add_shape(ctx.MSO_SHAPE.ROUNDED_RECTANGLE, left, top, column_width, height)
                card.fill.solid()
                card.fill.fore_color.rgb = card_fill
                card.line.fill.background()
                card.shadow.inherit = False
            if not items:
                continue
            pad = ctx.Inches(0.25)
            box = slide.shapes.add_textbox(left + pad, top + pad, column_width - (2 * pad), height - (2 * pad))
            self._add_colored_bullets(box.text_frame, items, ctx, size=16)
        return slide

    def _render_process_slide(self, prs, title_only_layout, title, body_texts, ctx):
        """Numbered step boxes arranged left to right instead of a flat
        bulleted list, per the Layout Classifier (ADR-023).

        ADR-061 — each step number now sits inside a solid accent-
        colored circle badge instead of floating as plain colored
        text — matches how every reference deck presents a numbered
        process (a badge, never a bare digit), and gives the eye a
        clear per-step anchor point in a row that otherwise has no
        visual separation between steps."""
        slide = prs.slides.add_slide(title_only_layout)
        self._apply_background(slide, ctx)
        self._add_corner_decoration(slide, ctx, prs, small=True)
        self._add_title_with_accent(slide, ctx, title)

        steps = body_texts[:5]  # beyond 5, a single row stops being readable
        slide_width, slide_height = prs.slide_width, prs.slide_height
        margin = ctx.Inches(0.4)
        box_width = (slide_width - (2 * margin)) // len(steps)
        number_top = int(slide_height * 0.34)
        badge_size = ctx.Inches(0.55)
        text_top = number_top + badge_size + ctx.Inches(0.15)
        text_height = ctx.Inches(2.3)

        for idx, step_text in enumerate(steps):
            left = margin + (idx * box_width)
            cleaned_text = PROCESS_BULLET_PATTERN.sub("", step_text).lstrip(",: ").strip()
            cleaned_text = cleaned_text[0].upper() + cleaned_text[1:] if cleaned_text else step_text

            badge_left = left + (box_width - badge_size) // 2
            badge = slide.shapes.add_shape(ctx.MSO_SHAPE.OVAL, badge_left, number_top, badge_size, badge_size)
            badge.fill.solid()
            badge.fill.fore_color.rgb = ctx.accent_color
            badge.line.fill.background()
            badge.shadow.inherit = False
            btf = badge.text_frame
            btf.word_wrap = False
            bp = btf.paragraphs[0]
            bp.text = str(idx + 1)
            bp.alignment = ctx.PP_ALIGN.CENTER
            # White text on a solid accent-colored badge — always high
            # contrast regardless of theme, unlike text-on-background
            # (which has to match the theme's own text color).
            ctx.style_run(bp, size=20, color=ctx.RGBColor(0xFF, 0xFF, 0xFF), bold=True)

            text_box = slide.shapes.add_textbox(left, text_top, box_width, text_height)
            tf = text_box.text_frame
            tf.word_wrap = True
            tp = tf.paragraphs[0]
            tp.text = cleaned_text
            tp.alignment = ctx.PP_ALIGN.CENTER
            ctx.style_run(tp, size=14, color=ctx.text_color)
        return slide


class _RenderContext:
    """Bundles the python-pptx types and theme values every renderer
    needs, plus a single helper for applying font family/size/color/
    bold consistently — this is what guarantees every layout looks
    like one coherent product instead of four independently-styled
    experiments (the actual root cause behind the "mediocre quality"
    feedback: nothing enforced consistency before this)."""

    def __init__(self, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_SHAPE,
                 title_color, accent_color, background_color, font_name,
                 text_color=None, media=None, corner_style="circle",
                 stat_chip=False, gradient_stops=None, accent_rgb=None, layout_style="default"):
        self.Inches = Inches
        self.Pt = Pt
        self.Emu = Emu
        self.RGBColor = RGBColor
        self.PP_ALIGN = PP_ALIGN
        self.MSO_SHAPE = MSO_SHAPE
        self.title_color = title_color
        self.accent_color = accent_color
        self.background_color = background_color
        self.text_color = text_color or title_color
        self.font_name = font_name
        self.media = media
        # ADR-059 — per-theme visual style controls, not just colors.
        # corner_style: "blob" (gradient), "circle" (plain accent oval,
        # the original/default look), or "none" (no corner decoration
        # at all — for minimal themes where restraint IS the style).
        # stat_chip: statistics-slide numbers get a colored background
        # box instead of plain text, when the theme calls for it.
        self.corner_style = corner_style
        self.stat_chip = stat_chip
        self.gradient_stops = gradient_stops
        self.accent_rgb = accent_rgb  # raw (r,g,b) tuple, for _tint() — RGBColor itself isn't blendable
        self.layout_style = layout_style
        # ADR-029: image_ids already used elsewhere in this deck — passed
        # to MediaPort.search_image(exclude_ids=...) so the same photo
        # never appears twice across one presentation's slides.
        self.used_image_ids: set = set()

    def style_run(self, paragraph, size=None, color=None, bold=None):
        """Applies the theme font family to a paragraph's font
        (covers the whole paragraph even before individual runs
        exist, since python-pptx creates the run from paragraph.text
        assignment) plus any explicitly requested size/color/bold."""
        paragraph.font.name = self.font_name
        if size is not None:
            paragraph.font.size = self.Pt(size)
        if color is not None:
            paragraph.font.color.rgb = color
        if bold is not None:
            paragraph.font.bold = bold

    def style_run_direct(self, run, size=None, color=None, bold=None):
        """ADR-061 — same as style_run, but for an individual Run
        object rather than a whole Paragraph. Needed wherever a single
        paragraph has to mix two different colors in one line (the
        colored bullet marker vs. the bullet text itself, in
        _add_colored_bullets) — paragraph.font only ever applies one
        color to the whole line, so those two runs each need their own
        call to this instead."""
        run.font.name = self.font_name
        if size is not None:
            run.font.size = self.Pt(size)
        if color is not None:
            run.font.color.rgb = color
        if bold is not None:
            run.font.bold = bold
