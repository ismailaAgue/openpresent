"""
Tests for the Media Port — Phase 3.5 Tier 2 (ADR-025), revised for the
multi-provider, scored, dedup-aware contract (ADR-029).
"""

import io
from backend.adapters.media.null_media_adapter import NullMediaAdapter
from backend.adapters.media.unsplash_adapter import UnsplashProvider
from backend.models.media import ImageResult

FAKE_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d494844520000000100000001080600000"
    "01f15c4890000000a49444154789c6360000002000100ffff03000006"
    "000557bf0e0000000049454e44ae426082"
)


class FakeHttpClient:
    def __init__(self, search_status=200, search_json=None, image_status=200, image_content=FAKE_PNG):
        self.search_status = search_status
        self.search_json = search_json if search_json is not None else {
            "results": [{"id": "abc123", "urls": {"small": "https://images.unsplash.com/fake.jpg"},
                         "description": "mountains at sunset"}]
        }
        self.image_status = image_status
        self.image_content = image_content
        self.calls = []

    def get(self, url, headers=None, timeout=10):
        self.calls.append(url)
        if "api.unsplash.com" in url:
            return {"status_code": self.search_status, "json": self.search_json, "content": b""}
        return {"status_code": self.image_status, "content": self.image_content, "json": None}


# -- NullMediaAdapter -----------------------------------------------------

def test_null_media_adapter_always_unavailable():
    assert NullMediaAdapter().is_available() is False


def test_null_media_adapter_always_returns_none():
    assert NullMediaAdapter().search_image("mountains") is None


# -- UnsplashProvider: two-phase candidate discovery + fetch --------------

def test_unsplash_unavailable_without_access_key():
    provider = UnsplashProvider(access_key="", http_client=FakeHttpClient())
    assert provider.is_available() is False
    assert provider.search_candidates("mountains") == []


def test_unsplash_available_with_access_key():
    assert UnsplashProvider(access_key="fake-key", http_client=FakeHttpClient()).is_available() is True


def test_unsplash_returns_scored_candidates():
    client = FakeHttpClient()
    provider = UnsplashProvider(access_key="fake-key", http_client=client)
    candidates = provider.search_candidates("mountains")
    assert len(candidates) == 1
    assert candidates[0].image_id == "unsplash:abc123"
    assert "mountains" in candidates[0].metadata_text


def test_unsplash_fetch_bytes_downloads_the_winning_candidate():
    client = FakeHttpClient()
    provider = UnsplashProvider(access_key="fake-key", http_client=client)
    candidate = provider.search_candidates("mountains")[0]
    result = provider.fetch_bytes(candidate)
    assert result == FAKE_PNG


def test_unsplash_returns_empty_on_search_failure():
    client = FakeHttpClient(search_status=401)
    provider = UnsplashProvider(access_key="fake-key", http_client=client)
    assert provider.search_candidates("mountains") == []


def test_unsplash_never_raises_on_broken_client():
    class BrokenClient:
        def get(self, *a, **k):
            raise ConnectionError("simulated network failure")

    provider = UnsplashProvider(access_key="fake-key", http_client=BrokenClient())
    assert provider.search_candidates("mountains") == []
    assert provider.fetch_bytes(
        type("C", (), {"fetch_url": "https://x"})()
    ) is None


# -- End-to-end: export pipeline embeds an image via the new interface ----

def test_export_embeds_image_when_media_adapter_available(monkeypatch):
    from backend.engines.generate import generate_presentation
    from backend.adapters import registry as reg
    from pptx import Presentation

    class FakeMediaAdapter:
        def is_available(self):
            return True

        def search_image(self, query, exclude_ids=None):
            return ImageResult(image_bytes=FAKE_PNG, image_id="fake:1",
                                provider="fake", relevance_score=1.0)

    monkeypatch.setattr(reg, "get_media_adapter", lambda: FakeMediaAdapter())

    source = (
        "**Team Overview**\n\n"
        "**Our Culture**\n"
        "We value collaboration and continuous learning across every team.\n"
    ).encode("utf-8")

    recipe, pptx_bytes = generate_presentation(file_bytes=source, filename="team.txt", export_format="pptx")
    prs = Presentation(io.BytesIO(pptx_bytes))

    title_slide = prs.slides[0]
    pictures = [s for s in title_slide.shapes if s.shape_type == 13]
    assert len(pictures) == 1


def test_export_unaffected_when_media_adapter_unavailable():
    """Confirms the default (NullMediaAdapter) path produces zero
    embedded images — the normal, current, $0-cost behavior."""
    from backend.engines.generate import generate_presentation
    from pptx import Presentation

    source = (
        "**Team Overview**\n\n"
        "**Our Culture**\n"
        "We value collaboration and continuous learning across every team.\n"
    ).encode("utf-8")

    recipe, pptx_bytes = generate_presentation(file_bytes=source, filename="team.txt", export_format="pptx")
    prs = Presentation(io.BytesIO(pptx_bytes))

    for slide in prs.slides:
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pictures) == 0


def test_export_never_embeds_the_same_image_twice_in_one_deck(monkeypatch):
    """ADR-029 duplicate prevention, verified end-to-end: a media
    adapter that always returns the SAME image_id should still only
    ever be embedded once per deck — the exporter should stop asking
    for more once it sees the id repeat, not silently duplicate it."""
    from backend.engines.generate import generate_presentation
    from backend.adapters import registry as reg
    from pptx import Presentation

    calls = {"exclude_ids_seen": []}

    class FakeMediaAdapter:
        def is_available(self):
            return True

        def search_image(self, query, exclude_ids=None):
            calls["exclude_ids_seen"].append(set(exclude_ids or set()))
            if exclude_ids and "fake:1" in exclude_ids:
                return None  # only one distinct image available
            return ImageResult(image_bytes=FAKE_PNG, image_id="fake:1",
                                provider="fake", relevance_score=1.0)

    monkeypatch.setattr(reg, "get_media_adapter", lambda: FakeMediaAdapter())

    source = (
        "**Team Overview**\n\n"
        "**First Section**\nSome content here about the first topic in detail.\n\n"
        "**Second Section**\nSome content here about the second topic in detail.\n\n"
        "**Third Section**\nSome content here about the third topic in detail.\n"
    ).encode("utf-8")

    recipe, pptx_bytes = generate_presentation(file_bytes=source, filename="team.txt", export_format="pptx")
    prs = Presentation(io.BytesIO(pptx_bytes))

    total_pictures = sum(
        1 for slide in prs.slides for s in slide.shapes if s.shape_type == 13
    )
    assert total_pictures <= 1  # never more than one embed of the same image
    # Confirms exclude_ids was actually threaded through on later calls.
    assert any(calls["exclude_ids_seen"])
