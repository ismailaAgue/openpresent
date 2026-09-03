"""Design-pass tests — ADR-061.

Covers the "every slide should carry a consistent visual identity, not
just the cover" fixes: colored bullet markers (replacing the default
black dot), a small corner decoration on every non-title slide, tinted
comparison-column cards, and circle badges on process-slide numbers.
A real off-canvas positioning bug was caught here by checking actual
rendered shape coordinates, not just that export() didn't raise — see
test_small_corner_decoration_stays_within_slide_bounds.
"""

import io
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu
from backend.adapters.export.pptx_adapter import PptxExportAdapter
from backend.models.recipe import Recipe, Outline, Slide, ContentBlock, BlockType, StructureSource, Theme


def make_recipe(theme_id="neutral", slides=None) -> Recipe:
    outline = Outline(structure_source=StructureSource.AI_GENERATED, slides=slides or [
        Slide(order=1, title="Cover", content_blocks=[]),
        Slide(order=2, title="Details", content_blocks=[
            ContentBlock(type=BlockType.BULLET, text="First point"),
            ContentBlock(type=BlockType.BULLET, text="Second point"),
        ]),
    ])
    return Recipe.new(project_id="p1", source_text="Topic: test", outline=outline,
                       theme=Theme(color_set_id=theme_id), audience_type="general", language="en")


def _render(recipe) -> Presentation:
    output = PptxExportAdapter().export(recipe)
    return Presentation(io.BytesIO(output))


def _oval_shapes(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    return [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and "OVAL" in str(s.auto_shape_type)]


def test_bullet_slide_paragraphs_have_no_native_bullet_element():
    """The colored marker run replaces the default dot — if buNone
    isn't actually present, LibreOffice/PowerPoint would show BOTH a
    native bullet and the accent-colored square, a real double-bullet
    regression this test exists specifically to catch."""
    prs = _render(make_recipe())
    body_placeholder = prs.slides[1].placeholders[1]
    for p in body_placeholder.text_frame.paragraphs:
        pPr = p._p.find(qn("a:pPr"))
        assert pPr is not None
        assert pPr.find(qn("a:buNone")) is not None


def test_bullet_slide_has_a_colored_marker_run_before_the_text():
    prs = _render(make_recipe())
    body_placeholder = prs.slides[1].placeholders[1]
    first_para = body_placeholder.text_frame.paragraphs[0]
    runs = first_para.runs
    assert len(runs) == 2
    assert runs[0].text.strip() == "■"
    assert runs[1].text == "First point"


def test_small_corner_decoration_appears_on_content_slides():
    """Every non-title slide now gets the small variant — previously
    only the title slide (slide 1) had any decoration at all."""
    prs = _render(make_recipe("bold_violet_stats"))
    assert len(_oval_shapes(prs.slides[1])) == 1


def test_small_corner_decoration_stays_within_slide_bounds():
    """Regression test for a real bug caught before shipping: the
    first version of the small-variant positioning hardcoded 16:9
    (13.333in) slide width while python-pptx's actual default
    Presentation() is 10in wide (4:3) — the decoration would have
    rendered miles off-canvas, invisible, on every real generated
    deck. A small intentional overhang (the shape "peeks in" from just
    past the edge — the same deliberate treatment the original large
    corner blob already uses) is fine; a multi-inch miss is the actual
    bug class this guards against."""
    prs = _render(make_recipe("bold_violet_stats"))
    slide = prs.slides[1]
    ovals = _oval_shapes(slide)
    assert len(ovals) == 1
    shape = ovals[0]
    tolerance = Emu(914400)  # 1 inch — generous vs. the ~0.12in intentional overhang
    assert -tolerance <= shape.left <= prs.slide_width
    assert -tolerance <= shape.top <= prs.slide_height
    assert shape.left + shape.width <= prs.slide_width + tolerance
    assert shape.top + shape.height <= prs.slide_height + tolerance


def test_minimal_mono_theme_still_has_no_decoration_on_content_slides():
    """corner_style='none' must still mean none — small=True doesn't
    override a theme's explicit choice to have no decoration at all."""
    prs = _render(make_recipe("minimal_mono"))
    assert len(_oval_shapes(prs.slides[1])) == 0


def test_comparison_slide_has_two_tinted_cards():
    slides = [
        Slide(order=1, title="Cover", content_blocks=[]),
        Slide(order=2, title="Us vs Them", content_blocks=[
            ContentBlock(type=BlockType.BULLET, text="Faster"),
            ContentBlock(type=BlockType.BULLET, text="Cheaper"),
            ContentBlock(type=BlockType.BULLET, text="Slower"),
            ContentBlock(type=BlockType.BULLET, text="Pricier"),
        ], layout_type="comparison"),
    ]
    prs = _render(make_recipe("gradient_violet", slides=slides))
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    rects = [s for s in prs.slides[1].shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and "ROUNDED_RECTANGLE" in str(s.auto_shape_type)]
    assert len(rects) == 2


def test_process_slide_badges_are_solid_circles_with_white_numbers():
    """Uses a theme with corner_style='circle' (not 'blob') so the
    small corner decoration and the step badges are unambiguous by
    shape count alone — a gradient-blob theme would also add one more
    oval (the corner decoration itself, correctly present per
    ADR-061), which a naive "count all ovals" assertion would
    conflate with the step badges."""
    slides = [
        Slide(order=1, title="Cover", content_blocks=[]),
        Slide(order=2, title="How It Works", content_blocks=[
            ContentBlock(type=BlockType.BULLET, text="First, sign up"),
            ContentBlock(type=BlockType.BULLET, text="Then connect your data"),
        ], layout_type="process"),
    ]
    prs = _render(make_recipe("bold_violet_stats", slides=slides))  # corner_style="circle"
    ovals = _oval_shapes(prs.slides[1])
    badges = [o for o in ovals if o.text_frame.text.strip() in ("1", "2")]
    assert len(badges) == 2
    # Solid accent fill (not the default theme placeholder fill) —
    # visually confirmed white-on-accent by rendering this to an
    # actual image (see ADR-061's Verification section); this just
    # checks the badge is genuinely a filled shape, not an empty outline.
    for badge in badges:
        assert badge.fill.type is not None
