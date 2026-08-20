"""
Tests for Phase 3.5 Step 4 (ADR-022): Layout Classifier and visual layouts.
"""

import pytest
from backend.layout.layout_classifier import classify_layout
from backend.models.recipe import Slide, ContentBlock, BlockType


def bullets(*texts):
    return [ContentBlock(type=BlockType.BULLET, text=t) for t in texts]


def test_statistics_slide_detected():
    slide = Slide(order=2, title="Revenue", content_blocks=bullets(
        "Total revenue reached $4.2M this quarter",
        "Growth of 12% quarter over quarter",
        "Northeast region outperformed by 18%",
    ))
    assert classify_layout(slide) == "statistics"


def test_plain_bullet_slide_not_misclassified_as_statistics():
    """A slide that mentions one number in passing should NOT trigger
    the statistics layout — only a clear majority of stat-like bullets
    should."""
    slide = Slide(order=2, title="Team Overview", content_blocks=bullets(
        "Our team grew from a handful of people to a much larger group over several years",
        "We work across multiple offices with people from over 12 different countries",
        "The culture emphasizes collaboration or maybe long term thinking about approach",
    ))
    assert classify_layout(slide) == "bullet_list"


def test_comparison_slide_detected_by_title():
    slide = Slide(order=2, title="Renewable vs Fossil Fuels", content_blocks=bullets(
        "Solar costs have dropped significantly", "Coal remains heavily subsidized",
    ))
    assert classify_layout(slide) == "comparison"


def test_comparison_versus_spelled_out_also_detected():
    slide = Slide(order=2, title="Cats versus Dogs", content_blocks=bullets("a", "b"))
    assert classify_layout(slide) == "comparison"


def test_default_bullet_list_for_normal_content():
    slide = Slide(order=2, title="Introduction", content_blocks=bullets(
        "This section introduces the main topic of the presentation.",
        "It sets up the context for what follows in later sections.",
    ))
    assert classify_layout(slide) == "bullet_list"


def test_comparison_takes_priority_over_statistics():
    """If a slide's title signals comparison AND its bullets happen to
    be stat-heavy, comparison wins — the title is a stronger, more
    deliberate signal than incidental bullet content."""
    slide = Slide(order=2, title="2020 vs 2024 Growth", content_blocks=bullets(
        "Revenue grew 40%", "Costs dropped 15%",
    ))
    assert classify_layout(slide) == "comparison"


def test_process_detected_by_title_keyword():
    slide = Slide(order=2, title="Project Timeline", content_blocks=bullets(
        "Kickoff meeting scheduled for week one",
        "Design phase runs through week four",
    ))
    assert classify_layout(slide) == "process"


def test_process_detected_by_sequential_language():
    slide = Slide(order=2, title="Onboarding", content_blocks=bullets(
        "First, the new hire completes paperwork.",
        "Then, they meet their team.",
        "Finally, they take ownership of a project.",
    ))
    assert classify_layout(slide) == "process"


def test_process_not_triggered_by_single_sequential_word():
    """One bullet that happens to start with 'Next' shouldn't be
    enough — require at least two sequential markers, same discipline
    as the statistics majority-threshold check."""
    slide = Slide(order=2, title="Overview", content_blocks=bullets(
        "Next quarter looks promising for growth in this specific market segment",
        "The team remains focused on long term product quality goals",
    ))
    assert classify_layout(slide) == "bullet_list"


def test_export_renders_process_slide_as_numbered_steps():
    from backend.engines.generate import generate_presentation
    from pptx import Presentation
    import io

    source = (
        "**Onboarding**\n\n"
        "**Onboarding Process**\n"
        "First, the new hire completes paperwork and account setup.\n"
        "Then, they meet their team and manager for an introduction.\n"
        "Next, they begin shadowing a colleague for two weeks.\n"
        "Finally, they take full ownership of their first project.\n"
    ).encode("utf-8")

    recipe, pptx_bytes = generate_presentation(file_bytes=source, filename="o.txt", export_format="pptx")
    prs = Presentation(io.BytesIO(pptx_bytes))

    process_slide = prs.slides[1]
    textboxes = [s for s in process_slide.shapes if not s.is_placeholder and s.has_text_frame and s.text_frame.text.strip()]
    # 4 steps, each with a number badge + text box = 8 textboxes
    assert len(textboxes) == 8
    number_boxes = [s for s in textboxes if s.text_frame.text.strip().isdigit()]
    assert sorted(b.text_frame.text.strip() for b in number_boxes) == ["1", "2", "3", "4"]
    # Sequential marker word should be stripped from the displayed text
    all_text = " ".join(s.text_frame.text for s in textboxes)
    assert "First," not in all_text
    assert "Finally," not in all_text


def test_export_renders_statistics_slide_as_separate_textboxes():
    """End-to-end: confirms the export adapter actually draws distinct,
    non-overlapping shapes for a statistics slide, not just that the
    layout_type field gets set."""
    from backend.engines.generate import generate_presentation
    from pptx import Presentation
    import io

    source = (
        "**Report**\n\n"
        "**Revenue**\n"
        "- Total revenue reached $4.2M this quarter\n"
        "- Growth of 12% quarter over quarter\n"
        "- Northeast region outperformed by 18%\n"
    ).encode("utf-8")

    recipe, pptx_bytes = generate_presentation(file_bytes=source, filename="r.txt", export_format="pptx")
    prs = Presentation(io.BytesIO(pptx_bytes))

    revenue_slide = prs.slides[1]  # slide 0 is the title slide
    textboxes = [s for s in revenue_slide.shapes if not s.is_placeholder and s.has_text_frame and s.text_frame.text.strip()]
    assert len(textboxes) == 3  # three distinct stat callouts, not one merged list
    left_positions = sorted(s.left for s in textboxes)
    assert left_positions[0] != left_positions[1] != left_positions[2]  # genuinely side by side


def test_export_renders_comparison_slide_as_two_columns():
    from backend.engines.generate import generate_presentation
    from pptx import Presentation
    import io

    source = (
        "**Renewable vs Fossil Fuels**\n\n"
        "**Renewable vs Fossil Energy Costs**\n"
        "- Solar costs have dropped significantly in recent years\n"
        "- Wind power is now cheaper than coal in most markets today\n"
        "- Fossil fuels still receive significant government subsidies\n"
        "- Coal plants face rising maintenance costs as infrastructure ages\n"
    ).encode("utf-8")

    recipe, pptx_bytes = generate_presentation(file_bytes=source, filename="c.txt", export_format="pptx")
    prs = Presentation(io.BytesIO(pptx_bytes))

    comparison_slide = prs.slides[1]
    textboxes = [s for s in comparison_slide.shapes if not s.is_placeholder and s.has_text_frame and s.text_frame.text.strip()]
    assert len(textboxes) == 2  # exactly two columns
    assert textboxes[0].left != textboxes[1].left  # genuinely side by side, not stacked
