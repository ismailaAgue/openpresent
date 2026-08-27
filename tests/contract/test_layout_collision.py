"""
Collision-detection tests — ADR-026.

These exist because a real, severe bug reached deployment: title text
and images overlapping on both title slides and content slides,
confirmed by measuring actual shape coordinates in real generated
files (not assumed from code review). This test suite makes that bug
class structurally impossible to reintroduce silently — every layout
that can contain both a title and an image is checked for real
bounding-box overlap, using an actual embedded image (not a stub),
the same way the original bug was only visible with a real image
present.
"""

import io
from PIL import Image
from pptx import Presentation
from backend.engines.generate import generate_presentation


def _real_placeholder_image(width=800, height=600, color=(120, 150, 180)) -> bytes:
    """A genuinely valid, decodable image — the original bug only
    appeared with real image bytes; a fake/corrupt stub wouldn't
    exercise the same code path faithfully."""
    img = Image.new("RGB", (width, height), color=color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _bounding_boxes_overlap(shape_a, shape_b) -> bool:
    """Standard axis-aligned rectangle intersection test."""
    a_left, a_top, a_right, a_bottom = shape_a.left, shape_a.top, shape_a.left + shape_a.width, shape_a.top + shape_a.height
    b_left, b_top, b_right, b_bottom = shape_b.left, shape_b.top, shape_b.left + shape_b.width, shape_b.top + shape_b.height
    return not (a_right <= b_left or b_right <= a_left or a_bottom <= b_top or b_bottom <= a_top)


def _text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def _picture_shapes(slide):
    return [s for s in slide.shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE


def _generate_with_fake_image(monkeypatch, source_text: str, filename: str = "doc.txt"):
    from backend.adapters import registry as reg
    from backend.models.media import ImageResult

    class FakeMediaAdapter:
        def is_available(self):
            return True

        def search_image(self, query, exclude_ids=None):
            return ImageResult(image_bytes=_real_placeholder_image(), image_id=f"fake:{query}",
                                provider="fake", relevance_score=1.0)

    monkeypatch.setattr(reg, "get_media_adapter", lambda: FakeMediaAdapter())
    recipe, pptx_bytes = generate_presentation(file_bytes=source_text.encode("utf-8"), filename=filename)
    return Presentation(io.BytesIO(pptx_bytes))


def test_title_slide_text_never_overlaps_image(monkeypatch):
    """Regression test for the confirmed P0 bug: the title slide's
    title text and its image were measured overlapping in real
    generated output (title vertical range [2130425, 3600450] vs
    image range [548640, 2990088] — real EMU values from the bug)."""
    prs = _generate_with_fake_image(monkeypatch, "**Quantum Entanglement: The Nonlocal Structure of Reality**\n\nSome content here.")
    title_slide = prs.slides[0]
    text_shapes = _text_shapes(title_slide)
    pictures = _picture_shapes(title_slide)
    assert len(pictures) == 1, "expected an image on the title slide for this test to be meaningful"

    for text_shape in text_shapes:
        for picture in pictures:
            assert not _bounding_boxes_overlap(text_shape, picture), (
                f"Title text {text_shape.text_frame.text!r} overlaps the image — "
                f"text=({text_shape.left},{text_shape.top},{text_shape.width},{text_shape.height}) "
                f"image=({picture.left},{picture.top},{picture.width},{picture.height})"
            )


def test_content_slide_title_never_overlaps_image(monkeypatch):
    """Regression test: a content slide's title was measured with an
    image's bounding box completely containing the title's bounding
    box in real generated output (title top=274638 fully inside image
    range [0, 1920240])."""
    prs = _generate_with_fake_image(
        monkeypatch,
        "**Quantum Entanglement: The Nonlocal Structure of Reality**\n\n"
        "**Fundamentals**\n"
        "Quantum mechanics presents a description of nature that differs radically from classical physics, "
        "introducing probability and superposition as fundamental features of reality itself.\n",
    )
    content_slide = prs.slides[1]
    pictures = _picture_shapes(content_slide)
    assert len(pictures) == 1, "expected an image on this content slide for this test to be meaningful"

    title_shape = content_slide.shapes.title
    for picture in pictures:
        assert not _bounding_boxes_overlap(title_shape, picture), (
            f"Content slide title overlaps the image — "
            f"title=({title_shape.left},{title_shape.top},{title_shape.width},{title_shape.height}) "
            f"image=({picture.left},{picture.top},{picture.width},{picture.height})"
        )


def test_content_slide_body_text_never_overlaps_image(monkeypatch):
    """The body text box must also not overlap the image — not just
    the title. This is the case that exposed the separate height=0
    placeholder-corruption bug found during investigation."""
    prs = _generate_with_fake_image(
        monkeypatch,
        "**Report**\n\n"
        "**Overview**\n"
        "This section explains the key findings from our recent analysis of market conditions.\n",
    )
    content_slide = prs.slides[1]
    pictures = _picture_shapes(content_slide)
    assert len(pictures) == 1

    body_texts = [s for s in _text_shapes(content_slide) if s is not content_slide.shapes.title]
    assert body_texts, "expected a body text box on this slide for the test to be meaningful"
    for body_shape in body_texts:
        # A body text box with zero height (the original corruption
        # bug) would trivially "not overlap" anything — explicitly
        # guard against that degenerate case being mistaken for a pass.
        assert body_shape.height > 0, "body text box has zero height — this is the placeholder-corruption bug, not a real pass"
        for picture in pictures:
            assert not _bounding_boxes_overlap(body_shape, picture)


def test_long_title_still_does_not_overlap_image(monkeypatch):
    """The original bug affected even short titles, but the critique
    specifically flagged long titles as a compounding risk — the
    'Atlantic Slave Trade' deck's very long title was the worst case
    observed. Test explicitly against a long title."""
    prs = _generate_with_fake_image(
        monkeypatch,
        "**The Economic Foundations of the Atlantic Slave Trade: A System of Capital and Cruelty**\n\n"
        "Historical background content goes here for context.",
    )
    title_slide = prs.slides[0]
    text_shapes = _text_shapes(title_slide)
    pictures = _picture_shapes(title_slide)
    assert len(pictures) == 1

    for text_shape in text_shapes:
        for picture in pictures:
            assert not _bounding_boxes_overlap(text_shape, picture)


def test_no_shape_extends_past_slide_boundaries(monkeypatch):
    """Additional geometry guard from the critique's P0 list: no
    content outside the slide."""
    prs = _generate_with_fake_image(
        monkeypatch,
        "**Team Culture**\n\n**Our Values**\nWe believe in collaboration and continuous learning across every team we build.",
    )
    for slide in prs.slides:
        for shape in slide.shapes:
            has_real_text = shape.has_text_frame and shape.text_frame.text.strip()
            is_picture = shape.shape_type == 13
            if not has_real_text and not is_picture:
                continue  # decorative shapes (accent bar, corner circle) are intentionally small/positioned off-canvas in one corner by design
            assert shape.left >= -100000, f"shape starts off the left edge: {shape.left}"
            assert shape.top >= -100000, f"shape starts off the top edge: {shape.top}"
            assert shape.left + shape.width <= prs.slide_width + 100000, "shape extends past the right edge"
            assert shape.top + shape.height <= prs.slide_height + 100000, "shape extends past the bottom edge"
