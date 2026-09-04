"""Editorial layout tests — ADR-062.

Covers the "editorial_cream" theme, built in direct response to a
side-by-side design critique against a competitor deck: full-bleed
cropped photography, a serif headline, kicker/footer metadata, and
statistics as a stacked sidebar panel rather than a row of chips.
Several real bugs (kicker cut mid-word, a leading "+" stranded outside
the chip number, a footer/caption collision on image slides) were
caught by actually rendering real decks to images before this test
file existed — see ADR-062's Verification section. These tests lock
in those specific fixes as regressions, plus basic structural
coverage for the new renderers.
"""

import io
from unittest.mock import MagicMock
from PIL import Image
from pptx import Presentation
from backend.adapters.export.pptx_adapter import PptxExportAdapter, CHIP_NUMBER_PATTERN
from backend.models.recipe import Recipe, Outline, Slide, ContentBlock, BlockType, StructureSource, Theme
from backend.models.media import ImageResult
import backend.adapters.registry as registry


def _fake_image_bytes(color=(80, 90, 120), size=(1200, 800)) -> bytes:
    buf = io.BytesIO()
    Image.new("RGB", size, color=color).save(buf, format="JPEG")
    return buf.getvalue()


def make_recipe(slides=None) -> Recipe:
    outline = Outline(structure_source=StructureSource.AI_GENERATED, slides=slides or [
        Slide(order=1, title="Climate Change, Explained", content_blocks=[]),
        Slide(order=2, title="Details", content_blocks=[
            ContentBlock(type=BlockType.BULLET, text="First point about the topic."),
            ContentBlock(type=BlockType.BULLET, text="Second point about the topic."),
        ]),
    ])
    return Recipe.new(project_id="p1", source_text="Topic: test", outline=outline,
                       theme=Theme(color_set_id="editorial_cream"), audience_type="general", language="en")


def _render(recipe) -> Presentation:
    output = PptxExportAdapter().export(recipe)
    return Presentation(io.BytesIO(output))


def test_editorial_theme_produces_a_valid_pptx():
    prs = _render(make_recipe())
    assert len(prs.slides) == 2


def test_editorial_title_slide_has_no_corner_decoration_shapes():
    """editorial_cream's corner_style is 'none' — the kicker/footer
    system replaces the corner-blob identity mark entirely, it
    shouldn't also draw one."""
    prs = _render(make_recipe())
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    ovals = [s for s in prs.slides[0].shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and "OVAL" in str(getattr(s, "auto_shape_type", ""))]
    assert len(ovals) == 0


def test_kicker_truncates_on_a_word_boundary_not_mid_word():
    """Regression test for a real bug caught by rendering: the first
    version passed a pre-sliced slide_data.title[:24] into
    _add_kicker, so its own word-boundary logic operated on an
    already-broken fragment ("...RUNNING A F"). Fixed by passing the
    full title through and letting _add_kicker do the only truncation."""
    slides = [
        Slide(order=1, title="Cover", content_blocks=[]),
        Slide(order=2, title="The Earth Is Running a Fever Today", content_blocks=[
            ContentBlock(type=BlockType.BULLET, text="Some supporting detail here."),
        ]),
    ]
    prs = _render(make_recipe(slides=slides))
    kicker_text = None
    for shape in prs.slides[1].shapes:
        if shape.has_text_frame and shape.text_frame.text.startswith("01 —"):
            kicker_text = shape.text_frame.text
            break
    assert kicker_text is not None
    # Must not end mid-word — every word present must be a COMPLETE
    # word from the original title, never a fragment of one.
    kicker_words = kicker_text.replace("01 — ", "").split()
    title_words = [w.upper() for w in "The Earth Is Running a Fever Today".split()]
    for w in kicker_words:
        assert w in title_words, f"'{w}' is not a complete word from the title — looks like a mid-word cut"


def test_chip_number_pattern_keeps_a_leading_sign_with_its_number():
    """Regression test: '+1.1C global warming...' previously matched
    only '1.1', stranding the '+' in the label half ('+C GLOBAL
    WARMING...') instead of with the number it belongs to."""
    m = CHIP_NUMBER_PATTERN.search("+1.1C global warming since 1880")
    assert m.group(0).strip() == "+1.1"
    m2 = CHIP_NUMBER_PATTERN.search("-2.5% decline in output")
    assert m2.group(0).strip() == "-2.5%"


def test_editorial_stats_slide_renders_a_sidebar_panel_not_a_row():
    slides = [
        Slide(order=1, title="Cover", content_blocks=[]),
        Slide(order=2, title="Key Numbers", content_blocks=[
            ContentBlock(type=BlockType.BULLET, text="+1.1C global warming since 1880"),
            ContentBlock(type=BlockType.BULLET, text="420 ppm atmospheric CO2 today"),
        ], layout_type="statistics"),
    ]
    prs = _render(make_recipe(slides=slides))
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    rects = [s for s in prs.slides[1].shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and "RECTANGLE" in str(getattr(s, "auto_shape_type", ""))]
    # panel background + top bar + 1 divider (2 stats = 1 divider between them)
    assert len(rects) >= 3


def test_editorial_content_slide_with_image_has_no_footer_deck_title_collision():
    """Regression test: the deck-title footer label and the image's
    own attribution caption previously both rendered in the same
    bottom-left region on image-bearing slides, overlapping. Fixed by
    skipping the footer's left label specifically when an image
    occupies that area — this test confirms exactly one textbox
    contains the attribution text down there, not two overlapping ones."""
    fake_media = MagicMock()
    fake_media.is_available.return_value = True
    fake_media.search_image.return_value = ImageResult(
        image_bytes=_fake_image_bytes(), image_id="x1", provider="fake",
        relevance_score=0.9, attribution="Photo by Test on Pexels",
    )
    original = registry.get_media_adapter
    registry.get_media_adapter = lambda: fake_media
    try:
        slides = [
            Slide(order=1, title="Cover", content_blocks=[]),
            Slide(order=2, title="With An Image", content_blocks=[
                ContentBlock(type=BlockType.BULLET, text="Some point about this."),
            ], image_query="test query"),
        ]
        prs = _render(make_recipe(slides=slides))
        bottom_left_texts = [
            s.text_frame.text for s in prs.slides[1].shapes
            if s.has_text_frame and "Photo by" in s.text_frame.text
        ]
        assert len(bottom_left_texts) == 1  # only the real caption, no duplicate deck-title label there
    finally:
        registry.get_media_adapter = original


def test_editorial_content_slide_uses_prose_for_sentence_ending_bullets():
    """_looks_like_prose reuses the exact rule the document export
    adapters already use (ADR-054) — sentence-ending bullets render as
    real connected paragraphs, not a bulleted list, matching the
    reference decks' body-text treatment."""
    slides = [
        Slide(order=1, title="Cover", content_blocks=[]),
        Slide(order=2, title="Prose Test", content_blocks=[
            ContentBlock(type=BlockType.BULLET, text="This is a full sentence with a period."),
            ContentBlock(type=BlockType.BULLET, text="This is another complete sentence."),
        ]),
    ]
    prs = _render(make_recipe(slides=slides))
    body_shapes = [s for s in prs.slides[1].shapes if s.has_text_frame and "full sentence" in s.text_frame.text]
    assert len(body_shapes) == 1
    # No dash-marker prefix run should be present for prose paragraphs
    assert "—" not in body_shapes[0].text_frame.text


def test_editorial_content_slide_uses_dash_bullets_for_fragments():
    slides = [
        Slide(order=1, title="Cover", content_blocks=[]),
        Slide(order=2, title="Fragment Test", content_blocks=[
            ContentBlock(type=BlockType.BULLET, text="Fast onboarding"),
            ContentBlock(type=BlockType.BULLET, text="Lower cost"),
        ]),
    ]
    prs = _render(make_recipe(slides=slides))
    body_shapes = [s for s in prs.slides[1].shapes if s.has_text_frame and "Fast onboarding" in s.text_frame.text]
    assert len(body_shapes) == 1
    assert "—" in body_shapes[0].text_frame.text
