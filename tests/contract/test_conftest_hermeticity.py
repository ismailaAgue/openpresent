"""
Regression tests for the root conftest fixture itself (ADR-037) — the
gap it closes was a real CI failure: `test_export_unaffected_when_
media_adapter_unavailable` passed locally (this sandbox has no network
route to wikimedia.org, so the real default provider's call silently
failed and looked like "no provider") but failed in GitHub Actions
CI (real network access, Wikimedia's call actually succeeded and
embedded a real photo) — the test's "no keys configured" assumption
had gone stale the moment Wikimedia/Wikipedia became always-on
defaults (ADR-029/030), and nothing enforced the assumption was still
true. These tests prove the fix directly rather than trusting that
individual test files happen to not need live network.
"""

from backend.adapters import registry


def test_media_adapter_defaults_to_null_within_any_normal_test():
    """Proves the autouse fixture actually intercepts the real
    registry function — this is the exact call that silently hit live
    Wikimedia in CI before ADR-037's fix."""
    adapter = registry.get_media_adapter()
    assert type(adapter).__name__ == "NullMediaAdapter"
    assert adapter.is_available() is False


def test_research_adapter_defaults_to_null_within_any_normal_test():
    """Same fix, the other provider that was always-on by default
    (Wikipedia) — same silent-live-call risk in any test that never
    explicitly mocked get_research_adapter."""
    adapter = registry.get_research_adapter()
    assert type(adapter).__name__ == "NullResearchAdapter"
    assert adapter.is_available() is False


def test_ai_adapter_defaults_to_null_within_any_normal_test():
    adapter = registry.get_ai_adapter()
    assert type(adapter).__name__ == "NullAdapter"
    assert adapter.is_available() is False


def test_individual_tests_can_still_override_the_hermetic_default(monkeypatch):
    """Confirms the conftest docstring's claim: a test's own
    monkeypatch.setattr call still wins over the autouse fixture's
    baseline — this is what every existing FakeMediaAdapter/
    FakeAIAdapter-style test in this suite depends on continuing to
    work exactly as before."""
    class FakeMediaAdapter:
        def is_available(self):
            return True

        def search_image(self, query, exclude_ids=None):
            return None

    monkeypatch.setattr(registry, "get_media_adapter", lambda: FakeMediaAdapter())
    adapter = registry.get_media_adapter()
    assert type(adapter).__name__ == "FakeMediaAdapter"
    assert adapter.is_available() is True


def test_generate_presentation_embeds_zero_images_with_real_default_registry_wiring():
    """End-to-end version of the same fix, going through the actual
    engine (not just checking the registry function directly) — this
    is effectively the original failing test, now passing for the
    correct reason (an explicitly-forced Null default) rather than by
    accident of this sandbox's own network restrictions."""
    import io
    from pptx import Presentation
    from backend.engines.generate import generate_presentation

    source = (
        "**Team Overview**\n\n"
        "**Our Culture**\n"
        "We value collaboration and continuous learning across every team.\n"
    ).encode("utf-8")

    recipe, pptx_bytes = generate_presentation(file_bytes=source, filename="team.txt", export_format="pptx")
    prs = Presentation(io.BytesIO(pptx_bytes))
    for slide in prs.slides:
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pictures) == 0
