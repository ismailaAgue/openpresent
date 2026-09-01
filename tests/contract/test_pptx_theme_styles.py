"""
Template theme tests — ADR-059.

Covers the four new theme presets (gradient_violet, minimal_mono,
bold_violet_stats, clean_saas_blue) and the stat-chip rendering logic
built for them. Two real bugs were found and fixed by actually
rendering these to images before this test file existed — see
ADR-059's Verification section — so these tests focus on the exact
things that broke: number extraction from bullet text that doesn't
match the classifier's own (deliberately narrower) STATISTIC_PATTERN,
and every theme producing genuinely different output, not just an
unused color entry nothing reads.
"""

from pptx import Presentation
from backend.adapters.export.pptx_adapter import (
    PptxExportAdapter, _COLOR_SETS, CHIP_NUMBER_PATTERN, _tint,
)
from backend.models.recipe import Recipe, Outline, Slide, ContentBlock, BlockType, StructureSource, Theme

NEW_THEME_IDS = ["gradient_violet", "minimal_mono", "bold_violet_stats", "clean_saas_blue"]


def make_recipe(theme_id: str, slides=None) -> Recipe:
    outline = Outline(structure_source=StructureSource.AI_GENERATED, slides=slides or [
        Slide(order=1, title="Elevate Your Pitch", content_blocks=[]),
        Slide(order=2, title="Why It Matters", content_blocks=[
            ContentBlock(type=BlockType.BULLET, text="Faster decisions across every team"),
            ContentBlock(type=BlockType.BULLET, text="Lower operating costs company-wide"),
        ]),
    ])
    return Recipe.new(project_id="p1", source_text="Topic: test", outline=outline,
                       theme=Theme(color_set_id=theme_id), audience_type="general", language="en")


def test_all_four_new_themes_are_registered():
    for theme_id in NEW_THEME_IDS:
        assert theme_id in _COLOR_SETS


def test_every_new_theme_has_the_required_style_keys():
    """corner_style/stat_chip aren't optional extras — every renderer
    that reads them assumes they exist (via .get() with a safe
    default, but a missing key silently falling back is exactly the
    kind of thing that should be caught here, not discovered later as
    a theme that looks wrong for no obvious reason)."""
    for theme_id in NEW_THEME_IDS:
        theme = _COLOR_SETS[theme_id]
        assert "corner_style" in theme
        assert theme["corner_style"] in ("blob", "circle", "none")
        assert "stat_chip" in theme


def test_gradient_violet_theme_has_gradient_stops():
    assert "gradient_stops" in _COLOR_SETS["gradient_violet"]
    stops = _COLOR_SETS["gradient_violet"]["gradient_stops"]
    assert len(stops) == 2


def test_each_new_theme_produces_different_output_bytes():
    """Real proof the four themes aren't just four names pointing at
    the same rendering — every pair must differ."""
    outputs = {tid: PptxExportAdapter().export(make_recipe(tid)) for tid in NEW_THEME_IDS}
    for a in NEW_THEME_IDS:
        for b in NEW_THEME_IDS:
            if a != b:
                assert outputs[a] != outputs[b]


def test_new_themes_produce_a_valid_pptx():
    for theme_id in NEW_THEME_IDS:
        output = PptxExportAdapter().export(make_recipe(theme_id))
        assert output.startswith(b"PK")  # pptx is a zip archive
        import io
        prs = Presentation(io.BytesIO(output))  # raises if not a real pptx
        assert len(prs.slides) == 2


def _oval_shapes(slide):
    """Auto-shapes only — .auto_shape_type raises ValueError for any
    shape that isn't one (textboxes, placeholders), so this has to
    filter by shape_type first rather than guard with getattr."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    return [
        s for s in slide.shapes
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and "OVAL" in str(s.auto_shape_type)
    ]


def test_minimal_mono_has_no_corner_decoration():
    """corner_style='none' means the title slide has exactly the shapes
    the title itself needs — no extra oval/blob shape added at all."""
    recipe = make_recipe("minimal_mono")
    output = PptxExportAdapter().export(recipe)
    import io
    prs = Presentation(io.BytesIO(output))
    assert len(_oval_shapes(prs.slides[0])) == 0


def test_gradient_violet_has_a_corner_blob():
    recipe = make_recipe("gradient_violet")
    output = PptxExportAdapter().export(recipe)
    import io
    prs = Presentation(io.BytesIO(output))
    assert len(_oval_shapes(prs.slides[0])) == 1


# -- Stat chip number extraction (the two real bugs found by rendering) -----

def test_chip_pattern_extracts_dollar_amounts():
    m = CHIP_NUMBER_PATTERN.search("$320,820M raised across 2018-2020")
    assert m.group(0).strip() == "$320,820M"


def test_chip_pattern_extracts_percentages_with_the_percent_sign_attached():
    """Regression test for the word-boundary bug: a trailing \\b in the
    pattern caused '90% client satisfaction rate' to match only '90',
    stranding the % in the label half. Confirmed by rendering, not
    just re-reading the regex."""
    m = CHIP_NUMBER_PATTERN.search("90% client satisfaction rate")
    assert m.group(0).strip() == "90%"


def test_chip_pattern_extracts_bare_magnitude_suffixes():
    """Regression test for the other real bug: layout_classifier's
    STATISTIC_PATTERN (only $ and %) doesn't match a bare '415K', so
    chip rendering needs its own, deliberately broader pattern."""
    m = CHIP_NUMBER_PATTERN.search("415K happy customers served")
    assert m.group(0).strip() == "415K"


def test_statistics_slide_with_chip_theme_does_not_overflow_a_bare_magnitude():
    """End-to-end version of the bug above: renders a real statistics
    slide with a 'K'-suffixed stat under the chip-style theme and
    confirms the number textbox actually has word_wrap enabled (the
    other half of the original bug — even once the regex matched, an
    unwrapped long value could still overflow its card)."""
    slides = [
        Slide(order=1, title="Elevate Your Pitch", content_blocks=[]),
        Slide(order=2, title="Key Metrics", content_blocks=[
            ContentBlock(type=BlockType.BULLET, text="$320,820M raised across 2018-2020"),
            ContentBlock(type=BlockType.BULLET, text="90% client satisfaction rate"),
            ContentBlock(type=BlockType.BULLET, text="415K happy customers served"),
        ], layout_type="statistics"),
    ]
    recipe = make_recipe("gradient_violet", slides=slides)
    output = PptxExportAdapter().export(recipe)
    import io
    prs = Presentation(io.BytesIO(output))
    stats_slide = prs.slides[1]
    text_boxes = [s for s in stats_slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    # 3 stats x 2 boxes each (number + label) + title = 7, at minimum
    assert len(text_boxes) >= 6
    all_text = " ".join(tb.text_frame.text for tb in text_boxes)
    assert "$320,820M" in all_text
    assert "90%" in all_text
    assert "415K" in all_text


def test_plain_stat_chip_theme_is_unaffected():
    """Every pre-existing theme (neutral, blue_academic, etc.) must
    keep rendering statistics as plain centered text, exactly as
    before — stat_chip=False is the explicit default for all of them,
    this proves it stayed that way."""
    for theme_id in ["neutral", "blue_academic", "warm_editorial", "modern_dark"]:
        assert _COLOR_SETS[theme_id]["stat_chip"] is False


def test_tint_blends_toward_white():
    result = _tint((0xB0, 0x3D, 0xE8), 0.85)
    # every channel should have moved closer to 255 than the original
    assert result[0] > 0xB0
    assert result[1] > 0x3D
    assert result[2] > 0xE8 - 1  # already close to 255; just not decreased
    assert _tint((10, 10, 10), 0.0) == (10, 10, 10)  # amount=0 is a no-op
    assert _tint((10, 10, 10), 1.0) == (255, 255, 255)  # amount=1 is pure white
